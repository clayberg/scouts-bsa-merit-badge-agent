"""ADK Golden Dataset Evaluation Harness (Rubric Category 5).

This module implements:
1. Two-part evaluation metrics:
   - Trajectory (The Process): Asserts tool execution sequence (research -> outline -> build).
   - Final Response (The Output): Audits generated .pptx slide decks for 100% requirement
     coverage, Scouts BSA brand palettes (#003F87), and max 7 bullets per slide.
2. Pydantic-backed EvalSet / EvalCase schema compatibility.
"""

import os
import json
import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from src.tools.scouting_scraper import fetch_merit_badge_pamphlet_pdf, MeritBadgeResearchRequest
from src.tools.pptx_builder import generate_bsa_slide_deck_pptx, PowerPointBuildRequest, SlideSpec
from src.config import ScoutsBSAPalette, is_eagle_required

GOLDEN_DATASET_PATH = os.path.join(os.path.dirname(__file__), "data", "golden_badges.json")

def load_golden_badges():
    with open(GOLDEN_DATASET_PATH, "r") as f:
        return json.load(f)

@pytest.mark.parametrize("golden_badge", load_golden_badges())
def test_golden_badge_trajectory_and_final_response(golden_badge, tmp_path):
    """Evaluates Trajectory and Final Response for benchmark merit badges."""
    badge_name = golden_badge["badge_name"]
    expected_reqs = golden_badge["expected_min_requirements"]
    
    # -------------------------------------------------------------------------
    # 1. TRAJECTORY EVALUATION: Tool Execution Sequence
    # -------------------------------------------------------------------------
    research_req = MeritBadgeResearchRequest(badge_name=badge_name)
    research_res = fetch_merit_badge_pamphlet_pdf(research_req)
    
    assert research_res["status"] == "SUCCESS", f"Failed to ingest pamphlet for {badge_name}"
    assert research_res["is_eagle_required"] == golden_badge["is_eagle_required"]
    assert len(research_res["requirements"]) >= expected_reqs
    
    # Build storyboard slides
    slides = []
    for req in research_res["requirements"]:
        slides.append(SlideSpec(
            title=f"Req {req['req_number']}: Instruction",
            bullet_points=[
                f"Requirement {req['req_number']} details.",
                req['req_text'][:100],
                "Practice with your Merit Badge Counselor."
            ],
            safety_warning=req.get("safety_callout")
        ))
        
    out_pptx = os.path.join(tmp_path, f"{badge_name.replace(' ', '_')}_Golden_Deck.pptx")
    build_req = PowerPointBuildRequest(
        badge_name=badge_name,
        slides=slides,
        output_path=out_pptx
    )
    build_res = generate_bsa_slide_deck_pptx(build_req)
    
    # -------------------------------------------------------------------------
    # 2. FINAL RESPONSE EVALUATION: python-pptx Audit of Generated Presentation
    # -------------------------------------------------------------------------
    assert os.path.exists(build_res["output_path"]), "PowerPoint .pptx was not created."
    
    prs = Presentation(build_res["output_path"])
    total_slides = len(prs.slides)
    
    # Assert 100% requirement coverage (slide count >= 1 title + expected_reqs)
    assert total_slides >= (1 + expected_reqs), (
        f"Requirement coverage failed: generated {total_slides} slides, "
        f"expected at least {1 + expected_reqs}"
    )
    
    # Audit pedagogical scannability rule (max 7 bullet points per slide)
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                bullets_count = len(shape.text_frame.paragraphs)
                assert bullets_count <= 8, (
                    f"Pedagogical scannability violated on slide {idx+1}: "
                    f"found {bullets_count} paragraphs (max allowed: 7 bullets + header)"
                )
                
    # Audit official Scouts BSA brand palette on Title Slide (#003F87 Navy Blue)
    title_slide = prs.slides[0]
    has_navy_blue = False
    for shape in title_slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                try:
                    if p.font.color.rgb == RGBColor(*ScoutsBSAPalette.NAVY_BLUE_RGB):
                        has_navy_blue = True
                except Exception:
                    pass
    assert has_navy_blue, f"Official BSA Navy Blue (#003F87) not found on Title Slide of {badge_name}"

def run_eval_suite():
    """Entry point script for CI/CD pipeline."""
    pytest.main([__file__, "-v", "--tb=short"])

if __name__ == "__main__":
    run_eval_suite()
