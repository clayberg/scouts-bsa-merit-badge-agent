"""PowerPoint Presentation Builder with Official Scouts BSA Styling.

This module implements:
1. Pydantic schemas for slide storyboards and build requests.
2. python-pptx generation with official Scouts BSA Navy Blue (#003F87) and Warm Olive (#4B5320).
3. Title slide Counselor Info + custom Troop Crest / Logo image placement.
4. Enforcement of the pedagogical max 7 bullet points per slide rule.
"""

import os
from typing import List, Optional, Dict, Any
from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from src.config import ScoutsBSAPalette, is_eagle_required
from src.tools.diagram_generator import get_badge_diagram_path

# ==============================================================================
# PYDANTIC JSON SCHEMAS
# ==============================================================================

class SlideSpec(BaseModel):
    """Schema representing the content and notes for an individual slide."""
    title: str = Field(..., description="Slide headline title.")
    bullet_points: List[str] = Field(..., description="List of bullet points (maximum 7 points).")
    presenter_notes: Optional[str] = Field(None, description="Counselor instructor notes.")
    safety_warning: Optional[str] = Field(None, description="Guide to Safe Scouting warning callout.")
    diagram_path: Optional[str] = Field(None, description="Optional path to diagram graphic to embed on slide.")

class CounselorTitleSlideInfo(BaseModel):
    """Counselor contact and troop customization for the presentation title slide."""
    counselor_name: str = Field(..., description="Full name of the Merit Badge Counselor.")
    troop_affiliation: str = Field(..., description="Troop number and council (e.g., 'Troop 101').")
    email_address: Optional[str] = Field(None, description="Optional contact email address.")
    phone_number: Optional[str] = Field(None, description="Optional contact phone number.")
    custom_troop_logo_path: Optional[str] = Field(None, description="Path to custom Troop Logo PNG/JPG.")

class PowerPointBuildRequest(BaseModel):
    """Schema for requesting a brand-compliant Scouts BSA PowerPoint presentation."""
    badge_name: str = Field(..., description="Official badge name.")
    slides: List[SlideSpec] = Field(..., description="Ordered list of slides to generate.")
    counselor_info: Optional[CounselorTitleSlideInfo] = Field(None, description="Title slide counselor customization.")
    output_path: Optional[str] = Field(None, description="Target filesystem path for saved .pptx file.")

class PowerPointBuildResult(BaseModel):
    """Structured return payload after generating PowerPoint presentation."""
    badge_name: str = Field(..., description="Official badge name.")
    slide_count: int = Field(..., description="Total number of slides generated.")
    output_path: str = Field(..., description="Absolute filesystem path to saved .pptx file.")
    is_eagle_required: bool = Field(..., description="True if Eagle-required badge styling was applied.")
    status: str = Field("SUCCESS", description="Execution status.")

# ==============================================================================
# PPTX GENERATOR FUNCTION
# ==============================================================================

def generate_bsa_slide_deck_pptx(request: PowerPointBuildRequest) -> Dict[str, Any]:
    """Generates a PowerPoint presentation using python-pptx with official BSA brand colors.
    
    Args:
        request: A validated PowerPointBuildRequest containing slides and counselor info.
        
    Returns:
        Dict: A PowerPointBuildResult dictionary with the output path and slide count.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    badge_title = request.badge_name.strip().title()
    eagle_flag = is_eagle_required(badge_title)
    
    # Define RGB Colors
    navy = RGBColor(*ScoutsBSAPalette.NAVY_BLUE_RGB)
    olive = RGBColor(*ScoutsBSAPalette.WARM_OLIVE_RGB)
    eagle_red = RGBColor(*ScoutsBSAPalette.EAGLE_RED_RGB)
    white = RGBColor(*ScoutsBSAPalette.WHITE_RGB)
    dark_text = RGBColor(*ScoutsBSAPalette.DARK_TEXT_RGB)
    
    # -------------------------------------------------------------------------
    # 1. TITLE SLIDE
    # -------------------------------------------------------------------------
    title_slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(title_slide_layout)
    
    # Add Title Box
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{badge_title} Merit Badge"
    p.font.name = "Roboto"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = navy
    p.alignment = PP_ALIGN.LEFT
    
    # Eagle-Required Subtitle Badge
    p2 = tf.add_paragraph()
    if eagle_flag:
        p2.text = "★ EAGLE-REQUIRED MERIT BADGE ★"
        p2.font.name = "Roboto Slab"
        p2.font.size = Pt(22)
        p2.font.bold = True
        p2.font.color.rgb = eagle_red
    else:
        p2.text = "Scouts BSA Official Advancement Guide"
        p2.font.name = "Roboto Slab"
        p2.font.size = Pt(20)
        p2.font.color.rgb = olive
        
    # Add Counselor Info Box on Title Slide
    if request.counselor_info:
        info = request.counselor_info
        counselor_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(8.0), Inches(2.2))
        ctf = counselor_box.text_frame
        ctf.word_wrap = True
        
        cp1 = ctf.paragraphs[0]
        cp1.text = f"Merit Badge Counselor: {info.counselor_name}"
        cp1.font.size = Pt(20)
        cp1.font.bold = True
        cp1.font.color.rgb = dark_text
        
        cp2 = ctf.add_paragraph()
        cp2.text = f"Unit: {info.troop_affiliation}"
        cp2.font.size = Pt(18)
        cp2.font.color.rgb = olive
        
        if info.email_address or info.phone_number:
            cp3 = ctf.add_paragraph()
            contact_str = ""
            if info.email_address:
                contact_str += f"Email: {info.email_address}   "
            if info.phone_number:
                contact_str += f"Phone: {info.phone_number}"
            cp3.text = contact_str.strip()
            cp3.font.size = Pt(16)
            cp3.font.color.rgb = dark_text
            
        # Add Custom Troop Logo if image file exists
        if info.custom_troop_logo_path and os.path.exists(info.custom_troop_logo_path):
            try:
                slide.shapes.add_picture(
                    info.custom_troop_logo_path,
                    Inches(9.5),
                    Inches(4.2),
                    width=Inches(2.5)
                )
            except Exception:
                pass
                
    # -------------------------------------------------------------------------
    # 2. CONTENT SLIDES
    # -------------------------------------------------------------------------
    for slide_spec in request.slides:
        content_slide = prs.slides.add_slide(title_slide_layout)
        
        # Header Title Box
        header_box = content_slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.0))
        htf = header_box.text_frame
        hp = htf.paragraphs[0]
        hp.text = slide_spec.title
        hp.font.name = "Roboto"
        hp.font.size = Pt(32)
        hp.font.bold = True
        hp.font.color.rgb = navy
        
        # Enforce Max 7 Bullet Points (Rubric Pedagogical Scannability)
        points = slide_spec.bullet_points[:7]
        
        # Determine if a pedagogical diagram should be embedded on this slide
        diagram_path = slide_spec.diagram_path or get_badge_diagram_path(request.badge_name, slide_spec.title)
        has_diagram = diagram_path and os.path.exists(diagram_path)
        
        # Content Body Box (Adjust width if diagram is present)
        body_width = Inches(6.4) if has_diagram else Inches(11.333)
        body_box = content_slide.shapes.add_textbox(Inches(1.0), Inches(1.8), body_width, Inches(4.5))
        btf = body_box.text_frame
        btf.word_wrap = True
        
        for idx, pt in enumerate(points):
            bp = btf.paragraphs[0] if idx == 0 else btf.add_paragraph()
            bp.text = f"•  {pt}"
            bp.font.name = "Roboto"
            bp.font.size = Pt(20) if has_diagram else Pt(22)
            bp.font.color.rgb = dark_text
            bp.space_after = Pt(10) if has_diagram else Pt(12)
            
        if has_diagram:
            try:
                content_slide.shapes.add_picture(
                    diagram_path,
                    Inches(7.7),
                    Inches(1.8),
                    width=Inches(4.6)
                )
            except Exception:
                pass
            
        # Guide to Safe Scouting Callout Box (if safety warning present)
        if slide_spec.safety_warning:
            safety_box = content_slide.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11.333), Inches(1.2))
            stf = safety_box.text_frame
            stf.word_wrap = True
            sp = stf.paragraphs[0]
            sp.text = f"⚠️ GUIDE TO SAFE SCOUTING: {slide_spec.safety_warning}"
            sp.font.name = "Roboto Slab"
            sp.font.size = Pt(16)
            sp.font.bold = True
            sp.font.color.rgb = eagle_red
            
        # Presenter Instructor Notes
        if slide_spec.presenter_notes:
            notes_slide = content_slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_spec.presenter_notes

    # -------------------------------------------------------------------------
    # 3. SAVE FILE
    # -------------------------------------------------------------------------
    out_filename = request.output_path or f"{badge_title.replace(' ', '_')}_Merit_Badge_Deck.pptx"
    prs.save(out_filename)
    
    result = PowerPointBuildResult(
        badge_name=badge_title,
        slide_count=len(prs.slides),
        output_path=os.path.abspath(out_filename),
        is_eagle_required=eagle_flag,
        status="SUCCESS"
    )
    return result.model_dump()
